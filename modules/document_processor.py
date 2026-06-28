"""Document Processing Module - Extract text from PDF, DOCX, PPTX, TXT."""
import os
import fitz
from docx import Document
from pptx import Presentation
import streamlit as st


class DocumentProcessor:
    supported_extensions = [".pdf", ".docx", ".pptx", ".txt"]

    def extract_text(self, file_path, file_extension):
        try:
            ext = file_extension.lower()
            if ext == ".pdf":
                return self._extract_pdf(file_path)
            elif ext == ".docx":
                return self._extract_docx(file_path)
            elif ext == ".pptx":
                return self._extract_pptx(file_path)
            elif ext == ".txt":
                return self._extract_txt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            st.error(f"Error extracting text: {str(e)}")
            return ""

    def _extract_pdf(self, file_path):
        text = ""
        with fitz.open(file_path) as pdf:
            for i, page in enumerate(pdf):
                text += f"\n--- Page {i + 1} ---\n"
                text += page.get_text()
        return text.strip()

    def _extract_docx(self, file_path):
        doc = Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
        for table in doc.tables:
            for row in table.rows:
                text += "\n" + " ".join([cell.text for cell in row.cells])
        return text.strip()

    def _extract_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {i} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text.strip()

    def _extract_txt(self, file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read()


def process_uploaded_file(uploaded_file, save_dir="data/uploads"):
    os.makedirs(save_dir, exist_ok=True)
    file_path = os.path.join(save_dir, uploaded_file.name)

    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

    processor = DocumentProcessor()
    ext = os.path.splitext(uploaded_file.name)[1].lower()

    if ext not in processor.supported_extensions:
        st.error(f"Unsupported file type: {ext}")
        return {}

    with st.spinner(f"Extracting text from {uploaded_file.name}..."):
        extracted_text = processor.extract_text(file_path, ext)

    if not extracted_text:
        st.warning(f"No text extracted from {uploaded_file.name}")
        return {}

    return {
        "filename": uploaded_file.name,
        "size_mb": round(os.path.getsize(file_path) / (1024 * 1024), 2),
        "extension": ext,
        "text": extracted_text,
        "word_count": len(extracted_text.split()),
        "char_count": len(extracted_text),
    }
